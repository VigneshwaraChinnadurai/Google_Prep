"""
Generate an editable PowerPoint presentation (.pptx) for the MBA Project.
Topic: Survey and Analysis of Quantum Processing Integration with Large Language Models (LLMs)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PATH = os.path.join(BASE_DIR, '..', 'Presentation.pptx')

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_HEADER = RGBColor(0x16, 0x21, 0x3E)
ACCENT = RGBColor(0xE9, 0x45, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)

def set_slide_bg(slide, color=BG_DARK):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title_text):
    """Add a colored title bar at the top."""
    # Header shape
    left = Inches(0.3)
    top = Inches(0.3)
    width = Inches(12.7)
    height = Inches(1.0)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_HEADER
    shape.line.color.rgb = RGBColor(0x0F, 0x34, 0x60)
    shape.line.width = Pt(2)
    
    # Title text
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = ACCENT

def add_content_box(slide, text_lines, top=Inches(1.6), left=Inches(0.8)):
    """Add content text box."""
    txBox = slide.shapes.add_textbox(left, top, Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if line.startswith('##'):
            # Subheading
            run = p.add_run()
            run.text = line[2:]
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = WHITE
            p.space_before = Pt(16)
        elif line.startswith('- '):
            # Bullet point
            run = p.add_run()
            run.text = line
            run.font.size = Pt(16)
            run.font.color.rgb = LIGHT_GRAY
            p.space_before = Pt(6)
            p.level = 1
        elif line == '':
            p.space_before = Pt(8)
        else:
            run = p.add_run()
            run.text = line
            run.font.size = Pt(16)
            run.font.color.rgb = WHITE
            p.space_before = Pt(4)

def add_footer(slide):
    """Add footer text."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(8), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Vigneshwara Chinnadurai | MBA Project | Manipal University Jaipur"
    run.font.size = Pt(10)
    run.font.color.rgb = MID_GRAY

# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
set_slide_bg(slide)

# Main title
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Survey and Analysis of"
run.font.size = Pt(24)
run.font.color.rgb = WHITE

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Quantum Processing Integration with"
run.font.size = Pt(32)
run.font.bold = True
run.font.color.rgb = ACCENT

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Large Language Models (LLMs)"
run.font.size = Pt(32)
run.font.bold = True
run.font.color.rgb = ACCENT

p = tf.add_paragraph()
p.space_before = Pt(40)
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "MBA Project Report"
run.font.size = Pt(20)
run.font.color.rgb = WHITE

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(20)
run = p.add_run()
run.text = "Vigneshwara Chinnadurai"
run.font.size = Pt(18)
run.font.color.rgb = LIGHT_GRAY

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Roll No: 2414504298 | Elective: Analytics & Data Science"
run.font.size = Pt(14)
run.font.color.rgb = MID_GRAY

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Guide: Mr. Govind"
run.font.size = Pt(14)
run.font.color.rgb = MID_GRAY

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(30)
run = p.add_run()
run.text = "Centre for Distance and Online Education | Manipal University Jaipur | May 2026"
run.font.size = Pt(12)
run.font.color.rgb = MID_GRAY

# ============================================================
# SLIDE 2: Introduction & Background
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, "Introduction & Background")
add_content_box(slide, [
    "##The Challenge:",
    "- LLMs (GPT-4, Gemini) require enormous computational resources",
    "- GPT-4 training estimated at $100M+ in compute costs",
    "- Exponential growth in model parameters (1.5B → 1.8T in 4 years)",
    "",
    "##The Opportunity:",
    "- Quantum computing offers fundamentally different computation",
    "- Superposition: n qubits represent 2ⁿ states simultaneously",
    "- Quantum parallelism could accelerate attention mechanisms",
    "- Potential for exponential speedup in specific sub-tasks",
    "",
    "##Research Question:",
    "- How can quantum computing be integrated with LLMs?",
    "- What is the current state of Quantum NLP research?",
])
add_footer(slide)

# ============================================================
# SLIDE 3: Research Objectives
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, "Research Objectives")
add_content_box(slide, [
    "",
    "- To comprehensively review and synthesize research on",
    "  quantum computing integration with LLMs (2017-2025)",
    "",
    "- To analyze and categorize existing approaches:",
    "  quantum-inspired, hybrid architectures, prototype QNLP models",
    "",
    "- To summarize technology trends and barriers to adoption",
    "",
    "- To conduct hands-on experimentation with quantum simulators",
    "  (IBM Qiskit, PennyLane, Google Cirq)",
    "",
    "- To provide strategic recommendations for future research",
    "  and practical integration in analytics/data science",
])
add_footer(slide)

# ============================================================
# SLIDE 4: Research Methodology
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, "Research Methodology")
add_content_box(slide, [
    "##Mixed-Methods Approach:",
    "",
    "##1. Systematic Literature Review",
    "- 47 papers analyzed from arXiv, IEEE, Google Scholar (2017-2025)",
    "- Thematic coding & technology maturity assessment",
    "",
    "##2. Experimental Research (4 Experiments)",
    "- Exp 1: Quantum Word Encoding (Amplitude, Angle, IQP)",
    "- Exp 2: Quantum Text Classification (Variational Circuits)",
    "- Exp 3: Hybrid Quantum-Classical Pipeline Comparison",
    "- Exp 4: Noise Analysis & Cross-Framework Benchmarking",
    "",
    "##Tools Used:",
    "- IBM Qiskit | PennyLane | Google Cirq | scikit-learn | Python",
])
add_footer(slide)

# ============================================================
# SLIDE 5: Literature Review Findings
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, "Literature Review: Key Findings")
add_content_box(slide, [
    "##Publication Growth:",
    "- 3 papers (2017-18) → 21 papers (2023-25): exponential growth",
    "",
    "##Research Distribution:",
    "- Theoretical/Framework: 29.8%",
    "- Simulation-Only: 38.3%",
    "- Hardware-Validated: 17.0%",
    "",
    "##Key Developments:",
    "- DisCoCat framework (Coecke et al., 2020): NLP → Quantum circuits",
    "- Quantum Transformers (Beer et al., 2021): Quantum attention",
    "- lambeq library: Practical QNLP toolkit",
    "- Quantinuum H2: First hardware-validated sentence classification",
])
add_footer(slide)

# ============================================================
# SLIDE 6: Experimental Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, "Key Experimental Results")
add_content_box(slide, [
    "##Experiment 1: Quantum Word Encoding",
    "- 94.2% encoding fidelity | 8.3:1 compression ratio (50d → 6 qubits)",
    "",
    "##Experiment 2: Quantum Text Classification",
    "- 87.3% accuracy with only 48 parameters (vs ~3000 for classical NN)",
    "- Competitive with SVM and Logistic Regression baselines",
    "",
    "##Experiment 3: Hybrid Quantum-Classical Pipeline",
    "- +8.5% accuracy advantage in low-data regime (n=50 samples)",
    "- 133x parameter reduction vs classical neural network",
    "",
    "##Experiment 4: Noise Resilience & Benchmarking",
    "- -5.8% accuracy drop at realistic noise level (p=0.01)",
    "- Hybrid approaches assessed at TRL 5-6 maturity",
])
add_footer(slide)

# ============================================================
# SLIDE 7: Conclusions
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, "Conclusions")
add_content_box(slide, [
    "",
    "##1. Quantum advantages are real but specific:",
    "- Encoding compression (94.2% fidelity, 8.3:1 ratio)",
    "- Small-data learning (+8.5% accuracy at n=50)",
    "- Parameter efficiency (133x reduction)",
    "",
    "##2. Hybrid approaches are the practical path forward",
    "- Best overall score in multi-criteria evaluation (27/35)",
    "- Competitive accuracy with dramatic parameter savings",
    "",
    "##3. Timeline: 10-15 years for full quantum LLMs",
    "- Near-term (1-3 yrs): Hybrid pilots on NISQ hardware",
    "- Medium-term (3-7 yrs): Quantum-enhanced sub-routines",
    "- Long-term (7-15 yrs): Fault-tolerant quantum transformers",
])
add_footer(slide)

# ============================================================
# SLIDE 8: Recommendations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title_bar(slide, "Recommendations")
add_content_box(slide, [
    "##For Organizations:",
    "- Establish quantum literacy programs for data science teams",
    "- Identify NLP use cases with small-data characteristics",
    "- Launch hybrid pilot projects using cloud quantum platforms",
    "- Develop 5-year quantum readiness roadmaps",
    "",
    "##For Researchers:",
    "- Develop standardized QNLP benchmarks",
    "- Focus on noise-resilient algorithms for NISQ hardware",
    "- Explore quantum advantages for low-resource languages",
    "",
    "##Phased Adoption Framework:",
    "- Phase 1 (2025-27): Literacy + Simulators",
    "- Phase 2 (2027-30): Hybrid Pilots on NISQ",
    "- Phase 3 (2030+): Production Quantum-Enhanced NLP",
])
add_footer(slide)

# ============================================================
# SLIDE 9: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Thank You"
run.font.size = Pt(48)
run.font.bold = True
run.font.color.rgb = ACCENT

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(20)
run = p.add_run()
run.text = "Questions & Discussion"
run.font.size = Pt(24)
run.font.color.rgb = WHITE

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(50)
run = p.add_run()
run.text = "Vigneshwara Chinnadurai"
run.font.size = Pt(18)
run.font.color.rgb = LIGHT_GRAY

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Roll No: 2414504298 | MBA - Analytics & Data Science"
run.font.size = Pt(14)
run.font.color.rgb = MID_GRAY

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(20)
run = p.add_run()
run.text = "Centre for Distance and Online Education | Manipal University Jaipur"
run.font.size = Pt(12)
run.font.color.rgb = MID_GRAY

# ============================================================
# SAVE
# ============================================================
prs.save(OUTPUT_PATH)
print(f"Presentation saved to: {OUTPUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
